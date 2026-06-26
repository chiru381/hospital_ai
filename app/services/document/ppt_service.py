from pptx import Presentation


def read_ppt(path):

    prs = Presentation(path)

    text = []

    for slide in prs.slides:

        for shape in slide.shapes:

            if hasattr(shape, "text"):

                if shape.text.strip():

                    text.append(shape.text)

    return "\n".join(text)